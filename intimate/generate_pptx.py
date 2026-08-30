"""
여성청결제 시장 분석 및 리뉴얼 전략 고도화 PPTX 발표 자료 생성 스크립트.

이 스크립트는 16:9 와이드 비율의 모던 비즈니스 스타일로 총 10장의 슬라이드를 구성하며,
상세페이지 전개 순서(Storyline) 비교 슬라이드와 2x2 포지셔닝 맵 그래픽 도형 드로잉 슬라이드를 포함하여
'intimate/report/market_analysis_and_strategy.pptx' 파일로 저장합니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_advanced_presentation():
    prs = Presentation()
    
    # 16:9 와이드 비율 설정 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 디자인 컬러 팔레트 정의
    NAVY_DARK = RGBColor(15, 23, 42)      # #0F172A (헤더 / 배경)
    NAVY_MAIN = RGBColor(30, 41, 59)      # #1E293B (제목 / 메인)
    TEAL_MAIN = RGBColor(13, 148, 136)    # #0D9488 (포인트 / 액센트)
    TEAL_LIGHT = RGBColor(204, 251, 241)  # #CCFBF1 (포인트 배경)
    BLUE_LIGHT = RGBColor(241, 245, 249)  # #F1F5F9 (카드 배경)
    BLUE_CARD = RGBColor(248, 250, 252)   # #F8FAFC (서브 카드)
    TEXT_DARK = RGBColor(51, 65, 85)      # #334155 (본문)
    TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B (서브 본문)
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)
    BORDER_DARK = RGBColor(148, 163, 184)
    ACCENT_CORAL = RGBColor(225, 29, 72)  # #E11D48 (경고/결핍 포인트)
    ACCENT_YELLOW = RGBColor(254, 240, 138) # #FEF08A (하이라이트)
    TARGET_BLUE = RGBColor(37, 99, 235)   # #2563EB (타깃 강조 블루)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # 공통 함수: 슬라이드 헤더 추가
    # -------------------------------------------------------------
    def add_header(slide, category, title, subtitle):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL_MAIN
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.45))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY_DARK
        
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # Slide 1: 표지 (Cover)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY_DARK
    bg1.line.color.rgb = NAVY_DARK
    
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.8), Inches(0.8), Inches(0.08))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = TEAL_MAIN
    bar1.line.color.rgb = TEAL_MAIN
    
    cover_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(4.5))
    tf1 = cover_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "여성청결제 시장 상세페이지 경쟁력 분석 및\n리뉴얼 전략 마스터 플랜"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(18)
    
    p2 = tf1.add_paragraph()
    p2.text = "18개 주요 브랜드 전수 벤치마킹 기반 · 2x2 포지셔닝 맵 · 전개 순서(Storyline) 비교 · '클린 사이언스' 리뉴얼 스토리보드"
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEAL_LIGHT
    p2.space_after = Pt(40)
    
    p3 = tf1.add_paragraph()
    p3.text = "기준 타깃 제품: 라엘(Rael) 천연 여성청결제  |  작성일자: 2026. 08. 30"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # Slide 2: Executive Summary
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Executive Summary", "여성청결제 시장 핵심 진단 및 3대 전환 전략 요약", "18개 주요 브랜드 전수 조사 결과, 상위권 브랜드는 '정량 임상과 독자 메커니즘'으로 결제를 이끌어냅니다.")
    
    cards_data2 = [
        ("01. 시장 완판 공식", "정량 임상 & 기전 규명", [
            "· '순함' 대신 '칸디다 99.9% 항균 / 99% 탈취' 수치 직격",
            "· 알칼리화와 유해균 점막 흡착 기전 과학적 폭로",
            "· 단순 추출물을 넘어선 독자 특허 바이옴 및 빙하수 베이스",
            "· [본품+리필] 파우치 기획으로 실속 가치 제공 (바솔)"
        ], BLUE_LIGHT),
        ("02. 기준 제품(라엘) 진단", "독보적 클린 vs 기능 결핍", [
            "· [강점] COSMOS 천연 인증 & 8가지 미니멀 성분의 높은 신뢰",
            "· [강점] 미국 아마존 1위 & 올리브영 생리대 팬덤 보유",
            "· [결핍] 항균/탈취/질염 개선 구체적 정량 데이터 부재",
            "· [한계] 공급자 중심 '설명문' 구조로 구매 전환율 저하"
        ], RGBColor(254, 242, 242)),
        ("03. Next 독점 포지셔닝", "클린 사이언스 펨테크", [
            "· 'COSMOS 천연 유기농' + 'D-만노스 유해균 부착 방어' 융합",
            "· 2대 질염균 99.9% 항균 성적서 & 5초 퀵 린스 실증",
            "· [본품 150ml + 친환경 리필 150ml] 기획으로 객단가 상승",
            "· 7일 무료체험 100% 환불 보장제로 구매 저항 제로화"
        ], TEAL_LIGHT)
    ]
    
    card_w = Inches(3.64)
    card_h = Inches(5.0)
    for i, (tag, title, bullets, bg_col) in enumerate(cards_data2):
        left = Inches(0.8 + i * 4.0)
        top = Inches(1.7)
        
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = bg_col
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide2.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), card_w - Inches(0.5), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEAL_MAIN if i != 1 else ACCENT_CORAL
        p.space_after = Pt(4)
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_DARK
        p_t.space_after = Pt(12)
        
        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = TEXT_DARK
            p_b.space_after = Pt(8)

    # -------------------------------------------------------------
    # Slide 3: [신규] 주요 브랜드 상세페이지 전개 순서(Storyline Flow) 비교
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Storyline Architecture", "주요 브랜드 상세페이지 메시지 전개 순서(Top-to-Bottom) 비교", "상위 브랜드는 [고민 자극 → 원인 규명 → 특허/기술 → 임상 성적서 → 번들]의 5단계 전환 퍼널을 엄격히 따릅니다.")
    
    # 6개 대표 브랜드의 스토리라인 비교 표
    storyline_headers = ["브랜드 / 제품", "1단계 (Hero 후킹)", "2단계 (문제 공감)", "3단계 (기술/포뮬러)", "4단계 (임상/안전성 증명)", "5단계 (가치 제안/락인)"]
    
    storyline_rows = [
        ("★ [기준] 라엘\n(천연 여성청결제)", "COSMOS 천연 인증\n단 8가지 전성분", "성분 개수에 대한 의문\n복잡한 성분의 자극 우려", "단 8가지 전성분 공개\n코코넛 유래 식물 세정", "피부 저자극 테스트 완료\n(정량 시험수치 미표기)", "150ml 단품\n청결티슈 묶음 기획"),
        ("이너시아\n(더퓨어 3X)", "KAIST 여성 과학자 R&D\n펨테크 오리진", "씻어도 반복되는 재발\n세균 점막 흡착 기전 규명", "3X 바이옴 (빙하수 100%\n+ 락토 + D-만노스)", "99% 악취 소취 성적서\n24종 알러젠 0% 불검출", "프리미엄 단품\n라보셀 생리대 크로스셀"),
        ("클리티\n(젤링워시)", "5대 국내외 특허 엠블럼\n8대 임상 완료 선언", "7가지 Y존 자가진단\n분비물/가려움/냄새 체크", "특허 5종 배합 (제주해수염\n+ TEFLOSE® 점막보호)", "칸디다 99.3%, 소취 99.5%\n즉각보습 93.09% 성적서", "300/500ml 대용량\n1+1 리필 번들 파격할인"),
        ("질경이\n(에코아 워시)", "6,600만 개 판매 신화\n단순 세정 脫피 선언", "Y존 복합 노화 지적\n(탄력 저하, 칙칙함, 건조)", "11개국 특허 바이옴-13\n하이드롤라이즈드 콜라겐", "3대 안티에이징 임상\n(보습 66%, 탄력 치밀도)", "N+N 대량 번들\n최대 50% 세트 할인"),
        ("아토팜\n(매터니티 케어)", "산부인과 테스트 완료\n'아주 좋음' 엠블럼", "임산부 분비물/자극 불안\n점막 잔여물 공포 환기", "독자 특허 AMPamide™\n신바이오틱스 EWG 그린", "7초 퀵 린스 잔여물 0%\n유해균 3종 99% 항균", "올리브영 단독 할인 기획\n1등 민감케어 신뢰 락인"),
        ("바솔\n(이너 밸런싱 폼)", "올리브영 W케어 1등\n산뜻 밸런스 선점", "생리 전후 불쾌취/찝찝함\n손 문지름 마찰 자극", "락토바실러스 발효용해물\n+ 호주산 티트리 + 시카", "칸디다균 99.9% 항균\n녹색소비자연대 안전적합", "[본품 150ml + 리필 150ml]\n올영 단독 더블 기획")
    ]
    
    col_widths3 = [Inches(1.9), Inches(2.2), Inches(2.4), Inches(2.5), Inches(2.5), Inches(2.2)]
    top_start = Inches(1.7)
    
    # 테이블 헤더 그리기
    left_acc = Inches(0.8)
    for j, (h_text, w) in enumerate(zip(storyline_headers, col_widths3)):
        h_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_acc, top_start, w, Inches(0.4))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = NAVY_MAIN
        h_box.line.color.rgb = NAVY_MAIN
        
        tb = slide3.shapes.add_textbox(left_acc, top_start + Inches(0.05), w, Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        left_acc += w
        
    # 테이블 행 그리기
    for i, row in enumerate(storyline_rows):
        is_target = (i == 0)
        row_top = top_start + Inches(0.42 + i * 0.8)
        left_acc = Inches(0.8)
        
        for j, (cell_text, w) in enumerate(zip(row, col_widths3)):
            c_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_acc, row_top, w, Inches(0.76))
            c_box.fill.solid()
            if is_target:
                c_box.fill.fore_color.rgb = RGBColor(238, 242, 255) if j == 0 else RGBColor(254, 242, 242)
            else:
                c_box.fill.fore_color.rgb = BLUE_CARD if i % 2 == 0 else WHITE
            c_box.line.color.rgb = BORDER_COLOR
            
            tb = slide3.shapes.add_textbox(left_acc + Inches(0.06), row_top + Inches(0.06), w - Inches(0.12), Inches(0.64))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(9.5)
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = TARGET_BLUE if is_target else NAVY_DARK
                p.alignment = PP_ALIGN.CENTER
            else:
                p.font.color.rgb = TEXT_DARK
                p.alignment = PP_ALIGN.LEFT
            left_acc += w

    # -------------------------------------------------------------
    # Slide 4: [직접 그래픽 드로잉] 여성청결제 시장 2x2 포지셔닝 맵
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "2x2 Positioning Map", "여성청결제 시장 2x2 포지셔닝 맵 및 White Space 시각화", "라엘은 'COSMOS 천연 유기농' 자산에 'D-만노스 유해균 부착 방어'를 결합한 [클린 사이언스 펨테크]로 이동해야 합니다.")
    
    # 2x2 매트릭스 그래픽 배경 영역 (중앙 배치)
    map_left = Inches(1.2)
    map_top = Inches(1.8)
    map_w = Inches(10.933)
    map_h = Inches(5.1)
    
    # 4개 사분면 배경 박스 생성
    quad_w = Inches(5.35)
    quad_h = Inches(2.4)
    gap_center = Inches(0.23)
    
    # Q1: 좌상단 (메디컬/고기능성 문제해결)
    q1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left, map_top, quad_w, quad_h)
    q1.fill.solid()
    q1.fill.fore_color.rgb = BLUE_LIGHT
    q1.line.color.rgb = BORDER_COLOR
    
    # Q2: 우상단 (White Space - 클린 사이언스 펨테크 / Next 라엘)
    q2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + quad_w + gap_center, map_top, quad_w, quad_h)
    q2.fill.solid()
    q2.fill.fore_color.rgb = TEAL_LIGHT
    q2.line.color.rgb = TEAL_MAIN
    q2.line.width = Pt(2)
    
    # Q3: 좌하단 (더마 / 약산성 마일드 케어)
    q3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left, map_top + quad_h + gap_center, quad_w, quad_h)
    q3.fill.solid()
    q3.fill.fore_color.rgb = BLUE_CARD
    q3.line.color.rgb = BORDER_COLOR
    
    # Q4: 우하단 (클린뷰티 오가닉 미니멀 / 현재 라엘)
    q4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + quad_w + gap_center, map_top + quad_h + gap_center, quad_w, quad_h)
    q4.fill.solid()
    q4.fill.fore_color.rgb = RGBColor(254, 242, 242)
    q4.line.color.rgb = ACCENT_CORAL
    
    # 축 레이블 텍스트 박스
    # Y축 상단 레이블
    yt_box = slide4.shapes.add_textbox(map_left + Inches(3.0), map_top - Inches(0.35), Inches(5.0), Inches(0.3))
    yt_p = yt_box.text_frame.paragraphs[0]
    yt_p.text = "▲ 문제성 집중 케어 (질염 / 칸디다 99.9% 항균 / 냄새 탈취 / 안티에이징)"
    yt_p.font.size = Pt(11)
    yt_p.font.bold = True
    yt_p.font.color.rgb = NAVY_DARK
    yt_p.alignment = PP_ALIGN.CENTER
    
    # Y축 하단 레이블
    yb_box = slide4.shapes.add_textbox(map_left + Inches(3.0), map_top + map_h + Inches(0.02), Inches(5.0), Inches(0.3))
    yb_p = yb_box.text_frame.paragraphs[0]
    yb_p.text = "▼ 데일리 마일드 케어 (0.00 저자극 / pH 약산성 / 온 가족 안심 세정)"
    yb_p.font.size = Pt(11)
    yb_p.font.bold = True
    yb_p.font.color.rgb = NAVY_DARK
    yb_p.alignment = PP_ALIGN.CENTER
    
    # X축 좌측 레이블
    xl_box = slide4.shapes.add_textbox(map_left - Inches(0.3), map_top + Inches(2.3), Inches(2.2), Inches(0.6))
    xl_p = xl_box.text_frame.paragraphs[0]
    xl_p.text = "◄ 더마/임상 과학\n   (Science & Medical)"
    xl_p.font.size = Pt(10.5)
    xl_p.font.bold = True
    xl_p.font.color.rgb = NAVY_DARK
    
    # X축 우측 레이블
    xr_box = slide4.shapes.add_textbox(map_left + map_w - Inches(1.8), map_top + Inches(2.3), Inches(2.2), Inches(0.6))
    xr_p = xr_box.text_frame.paragraphs[0]
    xr_p.text = "클린/비건/자연주의 ►\n(Clean & Natural)"
    xr_p.font.size = Pt(10.5)
    xr_p.font.bold = True
    xr_p.font.color.rgb = NAVY_DARK
    
    # 4개 사분면 내부 콘텐츠 및 브랜드 배지 배치
    # 1. 좌상단 (메디컬 문제해결)
    t1 = slide4.shapes.add_textbox(map_left + Inches(0.2), map_top + Inches(0.15), Inches(4.9), Inches(0.4))
    t1.text_frame.paragraphs[0].text = "클러스터 1: 메디컬 / 고기능 문제 해결형"
    t1.text_frame.paragraphs[0].font.size = Pt(12)
    t1.text_frame.paragraphs[0].font.bold = True
    t1.text_frame.paragraphs[0].font.color.rgb = NAVY_MAIN
    
    brands_q1 = ["질경이 (11개국 특허/안티에이징)", "이너시아 (KAIST D-만노스)", "클리티 (5대특허/8대임상)", "비레시피 (2대균 99.9% 항균)", "메디온 (락토리메디/질염케어)", "이너생각 (사상자/가려움 67%)"]
    for idx, b_name in enumerate(brands_q1):
        bx = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + Inches(0.2 + (idx%2)*2.5), map_top + Inches(0.65 + (idx//2)*0.52), Inches(2.4), Inches(0.42))
        bx.fill.solid()
        bx.fill.fore_color.rgb = WHITE
        bx.line.color.rgb = BORDER_DARK
        bx.text_frame.word_wrap = True
        bp = bx.text_frame.paragraphs[0]
        bp.text = b_name
        bp.font.size = Pt(9)
        bp.font.color.rgb = NAVY_DARK
        bp.alignment = PP_ALIGN.CENTER
        
    # 2. 우상단 (White Space - Next 라엘)
    t2 = slide4.shapes.add_textbox(map_left + quad_w + gap_center + Inches(0.2), map_top + Inches(0.15), Inches(4.9), Inches(0.4))
    t2.text_frame.paragraphs[0].text = "★ Next White Space: 클린 사이언스 펨테크"
    t2.text_frame.paragraphs[0].font.size = Pt(13)
    t2.text_frame.paragraphs[0].font.bold = True
    t2.text_frame.paragraphs[0].font.color.rgb = TEAL_MAIN
    
    target_next = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + quad_w + gap_center + Inches(0.4), map_top + Inches(0.65), Inches(4.5), Inches(1.5))
    target_next.fill.solid()
    target_next.fill.fore_color.rgb = WHITE
    target_next.line.color.rgb = TEAL_MAIN
    target_next.line.width = Pt(2)
    tf_tn = target_next.text_frame
    tf_tn.word_wrap = True
    ptn1 = tf_tn.paragraphs[0]
    ptn1.text = "★ [라엘(Rael) Next 목표 포지셔닝]"
    ptn1.font.size = Pt(12)
    ptn1.font.bold = True
    ptn1.font.color.rgb = TEAL_MAIN
    ptn2 = tf_tn.add_paragraph()
    ptn2.text = "· COSMOS 천연 유기농 무결점 안전성 +\n· 식물 유래 D-만노스 유해균 부착 방어 테크 +\n· 칸디다/가드넬라 99.9% 항균 성적서 & 24h 보습 결합"
    ptn2.font.size = Pt(10)
    ptn2.font.color.rgb = NAVY_DARK

    # 3. 좌하단 (더마 마일드)
    t3 = slide4.shapes.add_textbox(map_left + Inches(0.2), map_top + quad_h + gap_center + Inches(0.15), Inches(4.9), Inches(0.4))
    t3.text_frame.paragraphs[0].text = "클러스터 2: 더마 / 약산성 마일드 케어형"
    t3.text_frame.paragraphs[0].font.size = Pt(12)
    t3.text_frame.paragraphs[0].font.bold = True
    t3.text_frame.paragraphs[0].font.color.rgb = NAVY_MAIN
    
    brands_q3 = ["아토팜 (산부인과/7초 퀵린스)", "유리아쥬 (프랑스 온천수/만4세)", "일리윤 (공동특허 녹차락토스킨)", "해피바스 (8천원대 국민가성비)"]
    for idx, b_name in enumerate(brands_q3):
        bx = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + Inches(0.2 + (idx%2)*2.5), map_top + quad_h + gap_center + Inches(0.65 + (idx//2)*0.6), Inches(2.4), Inches(0.48))
        bx.fill.solid()
        bx.fill.fore_color.rgb = WHITE
        bx.line.color.rgb = BORDER_DARK
        bx.text_frame.word_wrap = True
        bp = bx.text_frame.paragraphs[0]
        bp.text = b_name
        bp.font.size = Pt(9.5)
        bp.font.color.rgb = NAVY_DARK
        bp.alignment = PP_ALIGN.CENTER

    # 4. 우하단 (클린뷰티 오가닉 미니멀 / 현재 라엘)
    t4 = slide4.shapes.add_textbox(map_left + quad_w + gap_center + Inches(0.2), map_top + quad_h + gap_center + Inches(0.15), Inches(4.9), Inches(0.4))
    t4.text_frame.paragraphs[0].text = "클러스터 3: 클린뷰티 오가닉 & 웰니스형"
    t4.text_frame.paragraphs[0].font.size = Pt(12)
    t4.text_frame.paragraphs[0].font.bold = True
    t4.text_frame.paragraphs[0].font.color.rgb = ACCENT_CORAL
    
    target_curr = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + quad_w + gap_center + Inches(0.2), map_top + quad_h + gap_center + Inches(0.6), Inches(2.5), Inches(1.55))
    target_curr.fill.solid()
    target_curr.fill.fore_color.rgb = WHITE
    target_curr.line.color.rgb = ACCENT_CORAL
    target_curr.line.width = Pt(1.5)
    tf_tc = target_curr.text_frame
    tf_tc.word_wrap = True
    ptc1 = tf_tc.paragraphs[0]
    ptc1.text = "🔴 [라엘 현재 위치]"
    ptc1.font.size = Pt(11)
    ptc1.font.bold = True
    ptc1.font.color.rgb = ACCENT_CORAL
    ptc2 = tf_tc.add_paragraph()
    ptc2.text = "· 단 8가지 성분\n· COSMOS 천연인증\n· [한계] 항균 수치 부재로\n  고민해결 고객 이탈"
    ptc2.font.size = Pt(9.5)
    ptc2.font.color.rgb = NAVY_DARK
    
    brands_q4 = ["바솔 (W케어 1등 티트리)", "아로마티카 (비건/에센셜)", "쏘피 (쿨링프레쉬 폼)", "디어스킨 (21개 성분 무향)"]
    for idx, b_name in enumerate(brands_q4):
        bx = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_left + quad_w + gap_center + Inches(2.85), map_top + quad_h + gap_center + Inches(0.6 + idx*0.4), Inches(2.2), Inches(0.35))
        bx.fill.solid()
        bx.fill.fore_color.rgb = WHITE
        bx.line.color.rgb = BORDER_DARK
        bx.text_frame.word_wrap = True
        bp = bx.text_frame.paragraphs[0]
        bp.text = b_name
        bp.font.size = Pt(9)
        bp.font.color.rgb = NAVY_DARK
        bp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # Slide 5: 주요 경쟁 그룹별 상세페이지 전략 심층 분석
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Competitor Analysis", "주요 3대 경쟁사 그룹별 상세페이지 소구 전략 비교", "고기능성 그룹은 임상 수치 폭격, 더마 그룹은 의학적 검증, 클린/매스 그룹은 감성과 가성비를 무기로 삼습니다.")
    
    group_data = [
        ("Group A. 고기능성 / 임상 소구 그룹", "대표 브랜드: 메디온, 질경이, 클리티, 비레시피, 이너시아", [
            "· [후킹] 1인칭 질염/가려움/냄새 고민 직격 (자가진단 체크리스트)",
            "· [원인] 알칼리화 및 세균 점막 흡착의 생물학적 기전 폭로",
            "· [증명] 칸디다 99.9% 항균 성적서 & 5대 등록 특허증 원본 노출",
            "· [전환] 100% 환불 보장제 및 1+1 더블 번들 세트 (객단가 2만원대)"
        ]),
        ("Group B. 마일드 / 더마 / 임산부 케어 그룹", "대표 브랜드: 아토팜, 유리아쥬, 일리윤, 디어스킨", [
            "· [후킹] 점막 자극 및 세정 후 건조함, 임산부 분비물 변화 공감",
            "· [성분] 프랑스 온천수, AMPamide™, 녹차 락토스킨 등 장벽 강화",
            "· [증명] 산부인과 사용 적합성 '아주 좋음' & 7초 퀵 린스 잔여물 제로",
            "· [전환] 300~500ml 메가 대용량 펌프로 온 가족 정착 유도"
        ]),
        ("Group C. 클린 / 자연유래 & 매스 그룹", "대표 브랜드: 아로마티카, 바솔, 쏘피, 해피바스", [
            "· [후킹] 인공 향료 불안, 생리 기간 꿉꿉함, 샤워 후 릴랙싱 니즈",
            "· [성분] 민들레 디콕션수, 티트리 오일, 쿨링 복합체, 국내산 쑥",
            "· [증명] 100% 비건 인증, 올리브영 W케어 1등 어워즈 엠블럼",
            "· [전환] 8~9천 원대 파격 최저가 or [본품+리필] 실속 기획 (바솔)"
        ])
    ]
    
    for i, (g_title, g_sub, g_bullets) in enumerate(group_data):
        top = Inches(1.8 + i * 1.7)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE_CARD
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide5.shapes.add_textbox(Inches(1.0), top + Inches(0.12), Inches(11.3), Inches(1.25))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = g_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_DARK
        
        p_sub = tf.add_paragraph()
        p_sub.text = g_sub
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = TEAL_MAIN
        p_sub.space_after = Pt(4)
        
        p_b = tf.add_paragraph()
        p_b.text = "   ".join(g_bullets[:2]) + "\n" + "   ".join(g_bullets[2:])
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # Slide 6: 핵심 6대 항목별 벤치마킹 매트릭스 요약
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Benchmarking Matrix", "핵심 6대 상세페이지 구성 요소별 상위권 성공 공식", "상위 브랜드들은 시각적 증명 장치와 다층 번들 전략으로 이탈을 방지하고 전환율을 극대화합니다.")
    
    matrix_items = [
        ("01. 후킹 & 헤드카피", "고민 직격 & 정량 수치형", "첫 3초 안에 '칸디다 99.9% 항균' 또는 '산부인과 테스트 완료' 엠블럼 노출"),
        ("02. 스토리라인 흐름", "5단계 논리 전환 깔때기", "고민 공감 ──► 원인 기전 규명 ──► 독자 바이옴 솔루션 ──► 임상 실증 ──► 번들 락인"),
        ("03. 성분 & 기술 소구", "정제수 대체 & 부착 방어", "캐나다 빙하수/온천수 베이스 + D-만노스 유해균 부착 차단 바이오 메커니즘 3D"),
        ("04. 임상 데이터 시각화", "공인 성적서 스캔본 노출", "시험기관 성적서 원본 + 비포/애프터 균 배양 샬레 비교 사진으로 신뢰도 200% 확보"),
        ("05. 가격 & 용량 전략", "친환경 리필 & 메가 대용량", "[본품 150ml + 리필 150ml] 기획(바솔) 또는 300~500ml 대용량으로 100ml 단가 인하"),
        ("06. 이탈 방지 장치", "100% 환불 보증제 & 리뷰", "7일 무료체험 100% 환불제(이너생각) + 올영 1위 엠블럼 + 고민별 찐후기 큐레이션")
    ]
    
    for i, (m_tag, m_title, m_desc) in enumerate(matrix_items):
        r = i // 2
        c = i % 2
        left = Inches(0.8 + c * 5.95)
        top = Inches(1.8 + r * 1.7)
        
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE_LIGHT
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), Inches(5.35), Inches(1.25))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = m_tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEAL_MAIN
        
        p_t = tf.add_paragraph()
        p_t.text = m_title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_DARK
        p_t.space_after = Pt(3)
        
        p_d = tf.add_paragraph()
        p_d.text = m_desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # Slide 7: 기준 제품의 결핍(Gap) 분석
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Gap Analysis", "기준 제품(라엘)의 4대 핵심 결핍 및 고객 구매 저항 진단", "안전성은 뛰어나지만, 고객의 고통을 해결해 주는 구체적 수치와 번들 가치가 결여되어 있습니다.")
    
    gap_items = [
        ("Gap 1. 정량 임상 수치 부재", "고객 구매 저항: '정말 냄새와 가려움이 사라질까?'", "· 피부 저자극 0.00 외 항균율/탈취율 수치 일체 미표기\n· '순함'에만 의존하여 질염/분비물 고민 고객의 장바구니 이탈 발생\n· [해결책] 칸디다균 99.9% 항균 & 악취 99% 소취 성적서 전면 노출", ACCENT_CORAL),
        ("Gap 2. 독자 바이옴 스토리 부재", "고객 구매 저항: '물에 계면활성제 탄 기본 워시 아닌가?'", "· 정제수, 글리세린 등 범용 원료의 단순 나열에 그침\n· 이너시아(D-만노스), 아토팜(AMPamide) 대비 독점 R&D 부재\n· [해결책] 유해균 부착 차단 D-만노스 바이오 메커니즘 3D 삽입", ACCENT_CORAL),
        ("Gap 3. 씻김성(Quick Rinse) 증명 부족", "고객 구매 저항: '거품이 미끈거리거나 남지 않을까?'", "· 단순 펌핑 거품 사진 1~2장으로 세정 쾌감 소구 부족\n· 점막 세정 잔여물에 대한 민감 고객의 심리적 불안 잔존\n· [해결책] 5초 퀵 린스 잔여물 제로 비커 비교 실험 영상/컷 배치", NAVY_MAIN),
        ("Gap 4. 친환경 리필 및 번들 혜택 부족", "고객 구매 저항: '150ml는 너무 빨리 써서 번거롭다'", "· 150ml 단품(12,900원) 위주로 운영되어 객단가 및 락인 한계\n· 플라스틱 절감 친환경 리필 파우치 옵션 부재\n· [해결책] [본품 150ml + 친환경 리필 150ml] 기획(19,900원) 런칭", NAVY_MAIN)
    ]
    
    for i, (g_tag, g_res, g_desc, col) in enumerate(gap_items):
        r = i // 2
        c = i % 2
        left = Inches(0.8 + c * 5.95)
        top = Inches(1.8 + r * 2.5)
        
        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(254, 242, 242) if col == ACCENT_CORAL else BLUE_CARD
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.35), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = g_tag
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_r = tf.add_paragraph()
        p_r.text = g_res
        p_r.font.size = Pt(11)
        p_r.font.bold = True
        p_r.font.color.rgb = NAVY_DARK
        p_r.space_after = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = g_desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # Slide 8: 도입부 카피라이팅 Before & After 개선안
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Copywriting Strategy", "도입부 카피라이팅 혁신: Before & After 3대 컨셉", "공급자 중심 스펙 나열에서 벗어나, 고객의 결핍을 정조준하는 3가지 카피 포트폴리오를 제안합니다.")
    
    b_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.95))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = RGBColor(254, 242, 242)
    b_box.line.color.rgb = BORDER_COLOR
    
    tb_b = slide8.shapes.add_textbox(Inches(1.0), Inches(1.88), Inches(11.3), Inches(0.8))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    p_b0 = tf_b.paragraphs[0]
    p_b0.text = "[Before 현재 카피]  \"단 8가지 전성분, COSMOS 천연 인증 순한 여성청결제\""
    p_b0.font.size = Pt(12)
    p_b0.font.bold = True
    p_b0.font.color.rgb = ACCENT_CORAL
    p_b1 = tf_b.add_paragraph()
    p_b1.text = "→ 이탈 원인: 제품 스펙의 단순 설명문으로 고객의 냄새/가려움 고통 공감 실패 및 스크롤 동기 부족"
    p_b1.font.size = Pt(10.5)
    p_b1.font.color.rgb = TEXT_MUTED
    
    after_cards = [
        ("컨셉 A. 직관적 고민 해결형", "\"씻어도 반복되는 찝찝함과 냄새, 단 8가지 순수한 성분으로 비워내다\"", "불필요한 화학 성분은 덜어내고, Y존을 괴롭히는 유해균과 악취만 99.9% 완벽 케어.\nCOSMOS 천연 거품으로 매일 산뜻하고 개운하게."),
        ("컨셉 B. 더마/임상 수치 증명형 (추천 ⭐)", "\"COSMOS 천연 인증에 칸디다균 99.9% 항균을 더하다 — 완벽한 증명\"", "식물 유래 D-만노스로 유해균 점막 부착을 차단하고, 피부 자극 0.00%로 완성한 클린 바이오 포뮬러.\n단 8가지 성분으로 입증한 99% 탈취력."),
        ("컨셉 C. 데일리 클린 밸런스형", "\"내 몸 가장 연약한 곳이니까 — 단 8가지 성분으로 지키는 마이크로바이옴\"", "흡수율 높은 점막을 위해 유해 화학 성분은 0%, 꼭 필요한 유익균 보호막만 정직하게 처방.\n매일 닿아도 자극 없는 천연 버블 웰니스.")
    ]
    
    for i, (c_tag, c_head, c_sub) in enumerate(after_cards):
        top = Inches(2.95 + i * 1.35)
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = TEAL_LIGHT if i == 1 else BLUE_LIGHT
        box.line.color.rgb = TEAL_MAIN if i == 1 else BORDER_COLOR
        
        tb = slide8.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = c_tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEAL_MAIN if i == 1 else NAVY_MAIN
        
        p_h = tf.add_paragraph()
        p_h.text = c_head
        p_h.font.size = Pt(13)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_DARK
        
        p_s = tf.add_paragraph()
        p_s.text = c_sub
        p_s.font.size = Pt(10)
        p_s.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # Slide 9: 상세페이지 리뉴얼 스토리보드 (Wireframe)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Renewal Storyboard", "상세페이지 7개 섹션 마스터 와이어프레임 구조", "고객의 구매 심리 흐름에 맞춰 [후킹 ──► 원인 ──► 기술 ──► 임상 ──► 제형 ──► 안전 ──► 번들]로 재배치합니다.")
    
    sections = [
        ("Section 1. 도입부 / 3초 후킹", "3대 엠블럼 + \"COSMOS 천연인증에 칸디다 99.9% 항균을 더하다\""),
        ("Section 2. 문제 제기 & 원인 규명", "3대 고민(냄새/가려움/분비물) + 높은 경피 흡수율 & 유해균 점막 흡착 기전"),
        ("Section 3. 독자 솔루션 & 기술", "단 8가지 전성분 인포그래픽 + D-만노스 유해균 부착 차단 3D 메커니즘"),
        ("Section 4. 압도적 임상 데이터 실증", "칸디다 99.99% & 가드넬라 99.8% 항균 성적서 + 99% 소취 그래프 + 24h 보습"),
        ("Section 5. 텍스처 & 5초 퀵 린스", "손 마찰 제로 버블 + 5초 만에 잔여물 없이 씻겨 내려가는 투명 비커 실험 컷"),
        ("Section 6. 무결점 클린 안전성", "COSMOS NATURAL 인증서 원본 + 피부 자극 0.00 + 알러젠 24종 불검출"),
        ("Section 7. 번들 오퍼 & 이탈 방지", "[본품 150ml + 친환경 리필 150ml] 기획(19,900원) + 7일 100% 환불 보증")
    ]
    
    for i, (s_title, s_desc) in enumerate(sections):
        top = Inches(1.8 + i * 0.72)
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE_CARD if i % 2 == 0 else BLUE_LIGHT
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide9.shapes.add_textbox(Inches(1.0), top + Inches(0.08), Inches(11.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = s_title + "  |  "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = TEAL_MAIN
        
        run = p.add_run()
        run.text = s_desc
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = NAVY_DARK

    # -------------------------------------------------------------
    # Slide 10: 결론 및 Next Step (A/B 테스트 및 런칭 전략)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Conclusion & Roadmap", "A/B 테스트 제언 및 최종 리뉴얼 실행 로드맵", "런칭 후 데이터 기반 최적화를 통해 올리브영 여성청결제 1등 브랜드로 도약합니다.")
    
    t_box1 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.75), Inches(2.3))
    t_box1.fill.solid()
    t_box1.fill.fore_color.rgb = TEAL_LIGHT
    t_box1.line.color.rgb = BORDER_COLOR
    
    tb1 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.35), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "🧪 Test 1. 도입부 헤드카피 A/B 테스트"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_DARK
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "· A안 (고민 해결형): \"씻어도 반복되는 찝찝함과 냄새...\"\n· B안 (임상 증명형 - 추천 ⭐): \"COSMOS 인증에 99.9% 항균을 더하다\"\n· [목표] 성분 중심 스마트 컨슈머의 구매 전환율(CVR) 25% 상승 검증"
    p1_sub.font.size = Pt(10.5)
    p1_sub.font.color.rgb = TEXT_DARK
    
    t_box2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.75), Inches(2.3))
    t_box2.fill.solid()
    t_box2.fill.fore_color.rgb = BLUE_LIGHT
    t_box2.line.color.rgb = BORDER_COLOR
    
    tb2 = slide10.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.35), Inches(2.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🧪 Test 2. 올리브영 메인 썸네일 A/B 테스트"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_DARK
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "· 1안 (기능/임상 중심): 투명 용기 + 버블 텍스처 + [COSMOS/99.9%] 엠블럼\n· 2안 (감성/웰니스 중심): 미니멀 욕실 라이프스타일 + [미국 1위] 엠블럼\n· [목표] 올리브영 카테고리 검색 결과 내 클릭률(CTR) 30% 개선"
    p2_sub.font.size = Pt(10.5)
    p2_sub.font.color.rgb = TEXT_DARK
    
    road_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.3))
    road_box.fill.solid()
    road_box.fill.fore_color.rgb = NAVY_DARK
    road_box.line.color.rgb = NAVY_DARK
    
    tb_r = slide10.shapes.add_textbox(Inches(1.0), Inches(4.55), Inches(11.3), Inches(2.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    pr = tf_r.paragraphs[0]
    pr.text = "🚀 라엘(Rael) 여성청결제 초격차 4단계 실행 로드맵"
    pr.font.size = Pt(14)
    pr.font.bold = True
    pr.font.color.rgb = TEAL_LIGHT
    pr.space_after = Pt(6)
    
    steps = [
        "1. [임상 데이터 확보] 칸디다/가드넬라 99.9% 항균 및 99% 소취 공인 성적서 발급",
        "2. [상세페이지 리뉴얼] 7개 섹션 와이어프레임 적용 & B안(임상 증명형) 카피라이팅 탑재",
        "3. [실속 번들 패키징] [본품 150ml + 친환경 리필 150ml] 올리브영 단독 기획 런칭",
        "4. [데이터 최적화] 썸네일 및 헤드카피 A/B 테스트를 통한 전환율(CVR) 30% 개선"
    ]
    for s in steps:
        p_s = tf_r.add_paragraph()
        p_s.text = s
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = WHITE
        p_s.space_after = Pt(2)
        
    output_path = os.path.join("intimate", "report", "market_analysis_and_strategy.pptx")
    prs.save(output_path)
    print(f"PPTX 발표 자료(10장 구성) 생성 성공: {output_path}")

if __name__ == "__main__":
    create_advanced_presentation()
